"""
Génère presentation_theorie.pptx
5 sections avec formules : État de l'art · Feature Engineering & Leakage ·
Régression · Classification · Clustering
Design identique au Beamer : bande grisclair + titre bleu + trait accent
Formules en Unicode math (éditables dans PowerPoint)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette Beamer ────────────────────────────────────────────
BLEU       = RGBColor(0,   70, 140)
BLEU_ACC   = RGBColor(0,  102, 204)
GRIS_CLAIR = RGBColor(240, 245, 250)
NOIR       = RGBColor(0,   0,   0)
BLANC      = RGBColor(255, 255, 255)
ROUGE      = RGBColor(180,  30,  30)
VERT       = RGBColor(0,  120,  60)

W = Inches(13.33)
H = Inches(7.5)
FIGS = ("/Users/ardachammahamatteguene/Desktop/"
        "Bureau - MacBook Pro de Ardacham/Memoire/latex/figures/")

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────
# Primitives
# ─────────────────────────────────────────────────────────────

def new_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background; bg.fill.solid(); bg.fill.fore_color.rgb = BLANC
    return sl

def rect(sl, x, y, w, h, fill):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(sl, x, y, w, h, text, size=14, bold=False, italic=False,
        color=NOIR, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def img(sl, path, x, y, w=None, h=None):
    if w and h: sl.shapes.add_picture(path, x, y, w, h)
    elif w:     sl.shapes.add_picture(path, x, y, width=w)
    elif h:     sl.shapes.add_picture(path, x, y, height=h)

def cap(sl, x, y, w, text):
    txt(sl, x, y, w, Inches(0.3), text, size=11,
        italic=True, color=RGBColor(100,100,100), align=PP_ALIGN.CENTER)

TITLE_H   = Inches(0.82)
HRULE_H   = Pt(2)
CONTENT_Y = TITLE_H + Inches(0.12)

def title_bar(sl, text, num):
    rect(sl, 0, 0, W, TITLE_H, GRIS_CLAIR)
    txt(sl, Inches(0.35), Inches(0.14), W - Inches(0.7), Inches(0.58),
        text, size=24, bold=True, color=BLEU)
    rect(sl, 0, TITLE_H, W, HRULE_H, BLEU_ACC)
    txt(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
        str(num), size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

def section_header(sl, num, letter, titre, couleur=BLEU):
    rect(sl, 0, 0, W, H, BLANC)
    rect(sl, 0, Inches(2.8), W, Pt(2), BLEU_ACC)
    rect(sl, 0, Inches(4.6), W, Pt(2), BLEU_ACC)
    rect(sl, Inches(4.4), Inches(2.9), Inches(4.6), Inches(1.6), GRIS_CLAIR)
    txt(sl, Inches(4.4), Inches(3.05), Inches(4.6), Inches(0.55),
        letter, size=32, bold=True, color=BLEU_ACC, align=PP_ALIGN.CENTER)
    txt(sl, Inches(4.4), Inches(3.6), Inches(4.6), Inches(0.7),
        titre, size=22, bold=True, color=BLEU, align=PP_ALIGN.CENTER)
    txt(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
        str(num), size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

def bullets(sl, x, y, w, h, items, title=None, t_size=14, b_size=13):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; first = True
    if title:
        p = tf.paragraphs[0]; first = False; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = title
        r.font.size = Pt(t_size); r.font.bold = True; r.font.color.rgb = BLEU
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = PP_ALIGN.LEFT; r = p.add_run()
        r.text = "–  " + item
        r.font.size = Pt(b_size); r.font.color.rgb = NOIR

def block(sl, x, y, w, h, title, items, t_col=BLEU, b_size=13):
    HDR = Inches(0.38)
    rect(sl, x, y, w, HDR, t_col)
    txt(sl, x + Inches(0.1), y + Inches(0.04), w - Inches(0.2), HDR,
        title, size=14, bold=True, color=BLANC)
    rect(sl, x, y + HDR, w, h - HDR, GRIS_CLAIR)
    tb = sl.shapes.add_textbox(x + Inches(0.12), y + HDR + Inches(0.05),
                                w - Inches(0.24), h - HDR - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = "–  " + item
        r.font.size = Pt(b_size); r.font.color.rgb = NOIR

def formula_box(sl, x, y, w, h, formula, explication=None):
    """
    Boîte formule : fond grisclair + bordure gauche bleuaccent (3pt)
    formule en bleu bold size 16 · explication en noir italic size 12
    """
    rect(sl, x, y, w, h, GRIS_CLAIR)
    rect(sl, x, y, Pt(4), h, BLEU_ACC)
    fh = Inches(0.42) if explication else h - Inches(0.08)
    txt(sl, x + Inches(0.18), y + Inches(0.06), w - Inches(0.26), fh,
        formula, size=15, bold=True, color=BLEU, align=PP_ALIGN.CENTER)
    if explication:
        txt(sl, x + Inches(0.18), y + fh + Inches(0.04),
            w - Inches(0.26), h - fh - Inches(0.1),
            explication, size=12, italic=True, color=RGBColor(70,70,70))

def add_table(sl, x, y, w, h, headers, rows, fs=12, cw=None):
    nc = len(headers)
    t  = sl.shapes.add_table(len(rows)+1, nc, x, y, w, h).table
    if cw:
        for i, c in enumerate(cw): t.columns[i].width = c
    for j, hdr in enumerate(headers):
        cell = t.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = BLEU
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = BLANC
    for i, row in enumerate(rows):
        bg = BLANC if i % 2 == 0 else GRIS_CLAIR
        for j, val in enumerate(row):
            cell = t.cell(i+1, j); cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fs); r.font.color.rgb = NOIR


# ══════════════════════════════════════════════════════════════
# DIAPO 1 — TITRE
# ══════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, Inches(2.3), W, Pt(2), BLEU_ACC)
rect(sl, 0, Inches(5.2), W, Pt(2), BLEU_ACC)
txt(sl, Inches(0.8), Inches(2.5), W - Inches(1.6), Inches(1.4),
    "Fondements Théoriques & Méthodes",
    size=36, bold=True, color=BLEU, align=PP_ALIGN.CENTER)
txt(sl, Inches(0.8), Inches(3.98), W - Inches(1.6), Inches(0.55),
    "État de l'art  ·  Feature Engineering  ·  Régression  ·  Classification  ·  Clustering",
    size=17, italic=True, color=BLEU_ACC, align=PP_ALIGN.CENTER)
txt(sl, Inches(0.8), Inches(5.35), W - Inches(1.6), Inches(0.38),
    "Ardacham Mahamat Teguene", size=16, bold=True, color=NOIR, align=PP_ALIGN.CENTER)
txt(sl, Inches(0.8), Inches(5.75), W - Inches(1.6), Inches(0.32),
    "Master 1 MIASHS  —  Université de Montpellier Paul Valéry  —  Juin 2026",
    size=14, color=RGBColor(80,80,80), align=PP_ALIGN.CENTER)
txt(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
    "1", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 2 — PLAN
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Plan", 2)
sections = [
    ("A", "État de l'art", "ML en santé, familles d'algorithmes, SHAP"),
    ("B", "Feature Engineering & Data Leakage", "Encodage, variables métier, détection du leakage"),
    ("C", "Régression", "6 modèles, formules Ridge/Lasso/RF/XGB, métriques R²/RMSE/MAE"),
    ("D", "Classification", "Précision/Rappel/F1/AUC, SMOTE, class_weight"),
    ("E", "Clustering", "K-Means, CAH Ward, silhouette, Davies-Bouldin"),
]
for i, (letter, titre, desc) in enumerate(sections):
    y = CONTENT_Y + Inches(0.18) + i * Inches(1.12)
    rect(sl, Inches(0.35), y, Inches(0.55), Inches(0.95), BLEU)
    txt(sl, Inches(0.35), y + Inches(0.2), Inches(0.55), Inches(0.55),
        letter, size=24, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    txt(sl, Inches(1.1), y + Inches(0.06), Inches(7), Inches(0.38),
        titre, size=17, bold=True, color=BLEU)
    txt(sl, Inches(1.1), y + Inches(0.48), Inches(11.5), Inches(0.38),
        desc, size=13, italic=True, color=RGBColor(80,80,80))


# ══════════════════════════════════════════════════════════════
# SECTION A — ÉTAT DE L'ART
# ══════════════════════════════════════════════════════════════

sl = new_slide()
section_header(sl, 3, "A", "État de l'art")

# ── A1 : ML en santé ──────────────────────────────────────────
sl = new_slide()
title_bar(sl, "A — État de l'art : ML en santé", 4)

bullets(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(6.3), Inches(2.5),
    ["Obermeyer & Emanuel (2016) : ML dépasse les méthodes statistiques classiques\n    pour anticiper les événements cliniques",
     "Rajpurkar et al. — Nature Medicine (2022) : panorama de l'état de l'art :\n    diagnostic assisté par IA, prédiction risque chronique, personnalisation",
     "Booth et al. (2012) : maladies chroniques issues de comportements\n    quotidiens mesurables → ML peut modéliser ces trajectoires"],
    title="Littérature de référence", b_size=13)

bullets(sl, Inches(0.35), CONTENT_Y + Inches(2.75), Inches(6.3), Inches(2.6),
    ["Prédiction du risque cardiovasculaire (cholestérol, HTA)",
     "Détection précoce du diabète de type 2",
     "Stratification de patients par profil de mode de vie",
     "Double enjeu : prédire ET expliquer (interprétabilité)"],
    title="Applications en santé numérique", b_size=13)

img(sl, FIGS+"machine_learning.png", Inches(6.8), CONTENT_Y + Inches(0.1), w=Inches(6.15))
cap(sl, Inches(6.8), CONTENT_Y + Inches(4.55), Inches(6.15),
    "Vue d'ensemble des familles d'algorithmes ML  (d'après Rajpurkar et al., 2022)")

# ── A2 : Comparatif algorithmes ────────────────────────────────
sl = new_slide()
title_bar(sl, "A — Comparatif des algorithmes utilisés", 5)

add_table(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
          W - Inches(0.7), Inches(4.05),
    ["Algorithme", "Famille", "Type", "Avantage clé", "Limite"],
    [["Régression linéaire",  "Linéaire",     "Supervisé",  "Interprétable, baseline",     "Relations linéaires seulement"],
     ["Ridge (L2)",           "Régularisée",  "Supervisé",  "Stabilise coeff. instables",  "Ne sélectionne pas les variables"],
     ["Lasso (L1)",           "Régularisée",  "Supervisé",  "Sélection automatique vars.",  "Élimine variables corrélées"],
     ["ElasticNet (L1+L2)",   "Régularisée",  "Supervisé",  "Combine Ridge + Lasso",       "2 hyperparamètres à tuner"],
     ["Random Forest",        "Ensembliste",  "Supervisé",  "Non-linéaire, parallèle",     "Peu interprétable"],
     ["XGBoost",              "Ensembliste",  "Supervisé",  "Corrige erreurs séquentiel.", "Risque surapprentissage"],
     ["Logistic Regression",  "Linéaire",     "Classif.",   "Interprétable, rapide",       "Frontière de décision linéaire"],
     ["K-Means",              "Partitionnel", "Non superv.",  "Rapide, O(n)",              "k fixé, sensible aux outliers"],
     ["CAH Ward",             "Hiérarchique", "Non superv.",  "Dendrogramme lisible",      "O(n²) → 5 000 individus max"]],
    fs=12, cw=[Inches(2.2), Inches(1.6), Inches(1.25), Inches(3.7), Inches(4.1)])

block(sl, Inches(0.35), CONTENT_Y + Inches(4.3), Inches(12.6), Inches(1.1),
    "Interprétabilité — SHAP (Lundberg & Lee, 2017)",
    ["Issu de la théorie des jeux coopératifs — valeur de Shapley par variable",
     "Donne la contribution marginale de chaque variable pour chaque prédiction individuelle",
     "Plus riche que feature importance globale : explication locale + direction (signe positif ou négatif)"])


# ══════════════════════════════════════════════════════════════
# SECTION B — FEATURE ENGINEERING & DATA LEAKAGE
# ══════════════════════════════════════════════════════════════

sl = new_slide()
section_header(sl, 6, "B", "Feature Engineering\n& Data Leakage")

# ── B1 : Feature Engineering ──────────────────────────────────
sl = new_slide()
title_bar(sl, "B — Feature Engineering : 15 → 16 variables explicatives", 7)

add_table(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
          Inches(6.5), Inches(2.4),
    ["Transformation", "Détail", "Motivation"],
    [["Suppression id",        "Identifiant sans pouvoir prédictif",    "Éviter le surapprentissage"],
     ["Encodage gender→0/1",   "female=0, male=1",                      "Modèles numériques requis"],
     ["hypertension (new)",    "1 si syst.≥140 ou diast.≥90  (ESC)",   "Seuil médical officiel"],
     ["bmi_cat (new)",         "0=sous-poids, 1=normal, 2=surpoids, 3=obèse  (OMS)", "Catégories cliniques"],
     ["lifestyle_score (new)", "score composite steps+sleep+water−smoker−alcool−bmi", "Indice synthétique"]],
    fs=12, cw=[Inches(1.9), Inches(2.6), Inches(2.0)])

bullets(sl, Inches(0.35), CONTENT_Y + Inches(2.65), Inches(6.5), Inches(1.4),
    ["z-score pour modèles linéaires :  x' = (x − μ) / σ",
     "min-max pour K-Means :  x' = (x − xmin) / (xmax − xmin)"],
    title="Mise à l'échelle", b_size=13)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(4.25), Inches(6.5), Inches(1.1),
    "lifestyle_score = steps + sleep + water − smoker − alcohol − bmi_z − calories_z",
    "Toutes les variables standardisées avant combinaison → contribution équitable")

# Côté droit : SHAP formule
block(sl, Inches(7.0), CONTENT_Y + Inches(0.1), Inches(5.95), Inches(2.0),
    "Variable hypercholesterol — rôle en clustering",
    ["Définie comme : cholesterol > 200 mg/dL  (NCEP)",
     "Conservée UNIQUEMENT pour le clustering (description des profils)",
     "Exclue de la régression : sinon data leakage →  voir slide suivante"])

bullets(sl, Inches(7.0), CONTENT_Y + Inches(2.25), Inches(5.95), Inches(3.2),
    ["Dataset initial : 15 variables",
     "− 1 variable supprimée (id)",
     "+ 3 variables créées (hypertension, bmi_cat, lifestyle_score)",
     "= 16 features explicatives finales",
     "hypercholesterol : exclue régression, incluse clustering",
     "Cibles : cholesterol (régression)  |  disease_risk (classification)"],
    title="Bilan des features", b_size=13)

# ── B2 : Data Leakage ─────────────────────────────────────────
sl = new_slide()
title_bar(sl, "B — Data Leakage : détection et correction", 8)

block(sl, Inches(0.35), CONTENT_Y + Inches(0.1), W - Inches(0.7), Inches(1.0),
    "Définition — Data Leakage  (Kaufman et al., 2012)",
    ["Une feature contient directement ou indirectement l'information de la variable cible",
     "→ le modèle 'voit' la réponse pendant l'entraînement → performances artificiellement gonflées",
     "→ sur nouvelles données (sans la variable leakée) : effondrement total des performances"],
    t_col=ROUGE)

y2 = CONTENT_Y + Inches(1.25)
rect(sl, Inches(0.35), y2, Inches(5.9), Inches(2.4), GRIS_CLAIR)
rect(sl, Inches(0.35), y2, Pt(4), Inches(2.4), ROUGE)
txt(sl, Inches(0.55), y2 + Inches(0.1), Inches(5.6), Inches(0.28),
    "Leakage détecté dans ce projet", size=13, bold=True, color=ROUGE)
txt(sl, Inches(0.55), y2 + Inches(0.42), Inches(5.6), Inches(0.35),
    "hypercholesterol = 1  si  cholesterol > 200 mg/dL", size=15, bold=True, color=BLEU)
txt(sl, Inches(0.55), y2 + Inches(0.82), Inches(5.6), Inches(1.45),
    "→ c'est une transformation directe de la variable cible (cholesterol)\n"
    "→ inclure cette variable revient à 'donner la réponse' au modèle\n"
    "→ corrélation hypercholesterol / cholesterol = 0.87\n"
    "→ R² artificiellement gonflé à 0.675", size=13, color=NOIR)

rect(sl, Inches(6.5), y2, Inches(6.45), Inches(2.4), GRIS_CLAIR)
rect(sl, Inches(6.5), y2, Pt(4), Inches(2.4), VERT)
txt(sl, Inches(6.7), y2 + Inches(0.1), Inches(6.1), Inches(0.28),
    "Correction appliquée", size=13, bold=True, color=VERT)
txt(sl, Inches(6.7), y2 + Inches(0.45), Inches(6.1), Inches(1.7),
    "Étape 1 — Détection : analyse de la définition de chaque variable\n"
    "  + vérification des corrélations feature/cible\n\n"
    "Étape 2 — Suppression : hypercholesterol retirée du jeu\n"
    "  d'entraînement en régression\n\n"
    "Résultat : R² 0.675  →  ≈ 0  (résultat honnête)", size=13, color=NOIR)

formula_box(sl, Inches(0.35), y2 + Inches(2.6), W - Inches(0.7), Inches(0.85),
    "R²  avec leakage = 0.675    →    R²  sans leakage ≈ 0.0001",
    "Preuve que 0.675 était artificiel — aucun signal réel dans les données")

block(sl, Inches(0.35), CONTENT_Y + Inches(4.7), W - Inches(0.7), Inches(1.0),
    "Règle générale — Prévenir le data leakage",
    ["Toujours auditer la définition de chaque feature avant modélisation",
     "Vérifier les corrélations feature/cible — une corrélation > 0.5 est suspecte sur données synthétiques",
     "Ne jamais calculer des features à partir de la variable cible (même indirectement)"])


# ══════════════════════════════════════════════════════════════
# SECTION C — RÉGRESSION
# ══════════════════════════════════════════════════════════════

sl = new_slide()
section_header(sl, 9, "C", "Régression")

# ── C1 : Métriques ────────────────────────────────────────────
sl = new_slide()
title_bar(sl, "C — Régression : métriques d'évaluation", 10)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(12.65), Inches(0.9),
    "R²  =  1  −  Σᵢ (yᵢ − ŷᵢ)²  /  Σᵢ (yᵢ − ȳ)²",
    "R²=1 : prédiction parfaite  |  R²=0 : modèle prédit la moyenne  |  R²<0 : pire que la moyenne")

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(1.15), Inches(6.1), Inches(0.85),
    "RMSE  =  √[ (1/n) × Σᵢ (yᵢ − ŷᵢ)² ]",
    "En mg/dL — pénalise les grandes erreurs (carré). Ici RMSE ≈ 43.33 mg/dL")

formula_box(sl, Inches(6.55), CONTENT_Y + Inches(1.15), Inches(6.45), Inches(0.85),
    "MAE  =  (1/n) × Σᵢ |yᵢ − ŷᵢ|",
    "Erreur absolue moyenne — plus robuste aux outliers. Ici MAE ≈ 37.59 mg/dL")

txt(sl, Inches(0.35), CONTENT_Y + Inches(2.2), W - Inches(0.7), Inches(0.28),
    "Notations :  yᵢ = valeur observée  |  ŷᵢ = valeur prédite  |  ȳ = moyenne observée  |  n = nb d'observations",
    size=12, italic=True, color=RGBColor(80,80,80))

add_table(sl, Inches(0.35), CONTENT_Y + Inches(2.65),
          W - Inches(0.7), Inches(1.45),
    ["Valeur de R²", "Interprétation", "Dans ce projet"],
    [["R² = 1",        "Prédiction parfaite — modèle explique 100% de la variance",    "—"],
     ["R² = 0.3–0.7",  "Signal détecté — modèle utile mais imprécis",                  "Attendu sur données réelles (cholestérol)"],
     ["R² ≈ 0",        "Modèle prédit uniquement la moyenne — aucun signal",            "Résultat obtenu : tous les modèles ≈ 0.0001"],
     ["R² < 0",        "Pire que la moyenne — surapprentissage du bruit",               "Random Forest non tuné : R² = −0.005"]],
    fs=12, cw=[Inches(1.85), Inches(6.4), Inches(4.4)])

block(sl, Inches(0.35), CONTENT_Y + Inches(4.35), W - Inches(0.7), Inches(1.05),
    "Protocole : Split 80/20 + K-Fold k=5 + RandomizedSearchCV n_iter=20",
    ["Split 80/20 : 80 000 train / 20 000 test  (random_state=42 pour reproductibilité)",
     "K-Fold k=5 : 5 itérations × 16 000 individus/fold — estimation robuste du score hors-échantillon",
     "RandomizedSearchCV : 20 tirages aléatoires d'hyperparamètres × k=5 = 100 entraînements vs GridSearch exhaustif"])

# ── C2 : Régression Linéaire + Ridge + Lasso ──────────────────
sl = new_slide()
title_bar(sl, "C — Régression linéaire, Ridge (L2) et Lasso (L1)", 11)

# Baseline
rect(sl, Inches(0.35), CONTENT_Y + Inches(0.1), W - Inches(0.7), Inches(0.28), GRIS_CLAIR)
txt(sl, Inches(0.5), CONTENT_Y + Inches(0.13), Inches(4), Inches(0.22),
    "Régression linéaire — baseline", size=14, bold=True, color=BLEU)
txt(sl, Inches(5.5), CONTENT_Y + Inches(0.13), Inches(7.3), Inches(0.22),
    "Si modèles plus complexes ≤ baseline → complexité non justifiée (rasoir d'Occam)",
    size=13, italic=True, color=RGBColor(80,80,80))

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(0.52), Inches(6.1), Inches(1.35),
    "ℒ_OLS = Σᵢ (yᵢ − ŷᵢ)²",
    "Minimise la somme des carrés des résidus (OLS).\n"
    "Hypothèse forte : relation linéaire entre features et cible.")

formula_box(sl, Inches(6.55), CONTENT_Y + Inches(0.52), Inches(6.45), Inches(1.35),
    "ŷ = w₀ + w₁x₁ + w₂x₂ + … + wₚxₚ",
    "wⱼ = coefficient associé à la variable xⱼ.\n"
    "Résultat : Lasso fixe les 16 coeff. / 16 à zéro → prédiction = ȳ = 224.3 mg/dL")

# Ridge
rect(sl, Inches(0.35), CONTENT_Y + Inches(2.05), W - Inches(0.7), Inches(0.28), GRIS_CLAIR)
txt(sl, Inches(0.5), CONTENT_Y + Inches(2.08), Inches(5), Inches(0.22),
    "Ridge — Régularisation L2  (Hoerl & Kennard, 1970)", size=14, bold=True, color=BLEU)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(2.47), Inches(6.1), Inches(1.4),
    "ℒ_Ridge = Σᵢ (yᵢ − ŷᵢ)²  +  λ × Σⱼ wⱼ²",
    "λ ≥ 0 : force de régularisation (tuné : λ = 7197)\n"
    "Réduit les wⱼ vers 0 SANS jamais les annuler complètement.")

formula_box(sl, Inches(6.55), CONTENT_Y + Inches(2.47), Inches(6.45), Inches(1.4),
    "wⱼ_Ridge = wⱼ_OLS / (1 + λ)",
    "Réduction proportionnelle à λ.\n"
    "Plus λ grand → coefficients plus proches de 0 mais jamais nuls.")

# Lasso
rect(sl, Inches(0.35), CONTENT_Y + Inches(4.05), W - Inches(0.7), Inches(0.28), GRIS_CLAIR)
txt(sl, Inches(0.5), CONTENT_Y + Inches(4.08), Inches(5), Inches(0.22),
    "Lasso — Régularisation L1  (Tibshirani, 1996)", size=14, bold=True, color=BLEU)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(4.47), Inches(6.1), Inches(1.4),
    "ℒ_Lasso = Σᵢ (yᵢ − ŷᵢ)²  +  λ × Σⱼ |wⱼ|",
    "Pénalité L1 (valeur absolue) → certains wⱼ = 0 exactement\n"
    "→ sélection automatique de variables. Ici : 16/16 coeff = 0.")

formula_box(sl, Inches(6.55), CONTENT_Y + Inches(4.47), Inches(6.45), Inches(1.4),
    "wⱼ_Lasso = signe(wⱼ_OLS) × max(|wⱼ_OLS| − λ, 0)",
    "Seuillage doux (soft thresholding) :\n"
    "si |wⱼ_OLS| < λ  →  wⱼ_Lasso = 0  (variable éliminée)")

# ── C3 : ElasticNet + RF + XGBoost ────────────────────────────
sl = new_slide()
title_bar(sl, "C — ElasticNet, Random Forest et XGBoost", 12)

# ElasticNet
rect(sl, Inches(0.35), CONTENT_Y + Inches(0.1), W - Inches(0.7), Inches(0.28), GRIS_CLAIR)
txt(sl, Inches(0.5), CONTENT_Y + Inches(0.13), Inches(4), Inches(0.22),
    "ElasticNet — L1 + L2  (Zou & Hastie, 2005)", size=14, bold=True, color=BLEU)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(0.52), W - Inches(0.7), Inches(1.0),
    "ℒ_EN = Σᵢ (yᵢ − ŷᵢ)²  +  λ [ ρ × Σⱼ|wⱼ|  +  (1−ρ)/2 × Σⱼ wⱼ² ]",
    "ρ ∈ [0,1] : mixing parameter.  ρ=1 → Lasso.  ρ=0 → Ridge.  Utile quand on ne sait pas lequel convient.")

# Random Forest
rect(sl, Inches(0.35), CONTENT_Y + Inches(1.7), W - Inches(0.7), Inches(0.28), GRIS_CLAIR)
txt(sl, Inches(0.5), CONTENT_Y + Inches(1.73), Inches(4), Inches(0.22),
    "Random Forest — Bagging d'arbres  (Breiman, 2001)", size=14, bold=True, color=BLEU)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(2.12), Inches(6.1), Inches(1.4),
    "ŷ_RF = (1/B) × Σ_{b=1}^{B}  Tᵦ(x)",
    "B arbres indépendants entraînés sur sous-échantillons aléatoires.\n"
    "Moyenne des prédictions → réduction de la variance.")

formula_box(sl, Inches(6.55), CONTENT_Y + Inches(2.12), Inches(6.45), Inches(1.4),
    "FI_Gini(j) = Σ_{t∈arbres} [p(t) × ΔGini(t,j)]",
    "Feature importance Gini : pondération par la fréquence du nœud.\n"
    "Attention : surévalue les variables à haute cardinalité.")

# XGBoost
rect(sl, Inches(0.35), CONTENT_Y + Inches(3.7), W - Inches(0.7), Inches(0.28), GRIS_CLAIR)
txt(sl, Inches(0.5), CONTENT_Y + Inches(3.73), Inches(4), Inches(0.22),
    "XGBoost — Gradient Boosting séquentiel  (Chen & Guestrin, 2016)", size=14, bold=True, color=BLEU)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(4.12), Inches(6.1), Inches(1.4),
    "ŷ⁽ᵗ⁾  =  ŷ⁽ᵗ⁻¹⁾  +  η × Tₜ(x)",
    "η ∈ (0,1] : learning rate. Tₜ : t-ième arbre.\n"
    "Chaque arbre corrige les erreurs du précédent (boosting).")

formula_box(sl, Inches(6.55), CONTENT_Y + Inches(4.12), Inches(6.45), Inches(1.4),
    "ℒ⁽ᵗ⁾ = Σᵢ ℓ(yᵢ, ŷᵢ⁽ᵗ⁾) + Ω(Tₜ)",
    "Ω = régularisation interne L1/L2 sur les feuilles.\n"
    "Moins sensible au surapprentissage que RF non tuné.")

# K-Fold
formula_box(sl, Inches(0.35), CONTENT_Y + Inches(5.7), W - Inches(0.7), Inches(0.7),
    "CV-score = (1/k) × Σ_{i=1}^{k} scoreᵢ     avec k = 5  →  16 000 individus/fold",
    "Score moyen sur les k folds — estimation stable du R² hors-échantillon")

# ── C4 : Résultats Régression ─────────────────────────────────
sl = new_slide()
title_bar(sl, "C — Résultats : régression (variable cible : cholesterol)", 13)

add_table(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
          Inches(9.3), Inches(2.65),
    ["Modèle", "λ optimal", "R²", "RMSE (mg/dL)", "MAE (mg/dL)", "CV-R²"],
    [["Linéaire (baseline)",  "—",     "0.0001",   "43.33", "37.59", "≈ 0"],
     ["Ridge  (L2)",          "7 197", "0.0001",   "43.33", "37.59", "≈ 0"],
     ["Lasso  (L1)",          "9.54",  "−0.000",   "43.33", "37.60", "≈ 0"],
     ["ElasticNet",           "—",     "−0.000",   "43.33", "37.60", "≈ 0"],
     ["Random Forest (tuné)", "—",     "−0.005",   "43.44", "37.67", "≈ 0"],
     ["XGBoost (tuné)",       "—",     "0.0001",   "43.33", "37.59", "≈ 0"]],
    fs=13, cw=[Inches(2.5), Inches(1.1), Inches(0.95), Inches(1.55), Inches(1.55), Inches(1.55)])

block(sl, Inches(9.6), CONTENT_Y + Inches(0.1), Inches(3.38), Inches(2.65),
    "Lasso — résultat formel",
    ["λ=9.54 → 16/16 coeff = 0",
     "Prédit ȳ = 224.3 mg/dL pour tous",
     "Preuve : aucune variable\nlinéairement prédictive",
     "Sélection automatique = 0 variables"])

block(sl, Inches(0.35), CONTENT_Y + Inches(2.85), W - Inches(0.7), Inches(1.1),
    "Analyse critique",
    ["Écart entre R²=−0.005 (RF) et R²=0.0001 (XGBoost) : statistiquement négligeable → pas de meilleur modèle",
     "RF non tuné : R²<0 → surapprentissage du bruit. Tuning ramène à 0.0001 sans changer la conclusion",
     "Cause fondamentale : variables synthétiques générées indépendamment → aucun signal à modéliser"])

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(4.15), W - Inches(0.7), Inches(1.0),
    "R²_leakage = 0.675   →   R²_corrigé ≈ 0.0001   (après suppression de hypercholesterol)",
    "Leakage détecté en analysant les corrélations features/cible et la définition de la variable.")

img(sl, FIGS+"correlation_matrix.png",
    Inches(0.35), CONTENT_Y + Inches(5.35), w=Inches(12.65))
cap(sl, Inches(0.35), H - Inches(0.4), Inches(12.65),
    "Matrice de corrélations — toutes valeurs < |0.01|  →  confirme l'absence de signal linéaire")


# ══════════════════════════════════════════════════════════════
# SECTION D — CLASSIFICATION
# ══════════════════════════════════════════════════════════════

sl = new_slide()
section_header(sl, 14, "D", "Classification")

# ── D1 : Métriques ────────────────────────────────────────────
sl = new_slide()
title_bar(sl, "D — Classification : métriques d'évaluation", 15)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(4.0), Inches(1.45),
    "Précision = TP / (TP + FP)",
    "Parmi les individus classés 'à risque',\ncombien le sont réellement ?\n→ qualité des positifs prédits")

formula_box(sl, Inches(4.55), CONTENT_Y + Inches(0.1), Inches(4.0), Inches(1.45),
    "Rappel = TP / (TP + FN)",
    "Parmi les individus réellement à risque,\ncombien sont détectés ?\n→ capacité à ne pas rater un cas")

formula_box(sl, Inches(8.75), CONTENT_Y + Inches(0.1), Inches(4.2), Inches(1.45),
    "F1 = 2 × Prec × Rappel\n      /(Prec + Rappel)",
    "Moyenne harmonique Précision / Rappel\n→ métrique équilibrée sur données déséquilibrées")

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(1.7), Inches(12.65), Inches(0.88),
    "AUC-ROC = ∫₀¹ TPR(FPR) dFPR     (aire sous la courbe ROC)",
    "AUC=1 : discrimination parfaite  |  AUC=0.5 : équivalent à un tirage aléatoire  |  AUC<0.5 : pire que le hasard")

txt(sl, Inches(0.35), CONTENT_Y + Inches(2.75), W - Inches(0.7), Inches(0.28),
    "TP = vrais positifs  |  FP = faux positifs  |  FN = faux négatifs  |  TN = vrais négatifs  |  TPR = TP/(TP+FN)  |  FPR = FP/(FP+TN)",
    size=12, italic=True, color=RGBColor(80,80,80))

add_table(sl, Inches(0.35), CONTENT_Y + Inches(3.1),
          W - Inches(0.7), Inches(1.25),
    ["Métrique", "Naïf (tout=0)", "Log. Regr.", "Random Forest", "XGBoost", "Interprétation"],
    [["Accuracy",    "0.752", "0.498", "0.568", "0.525", "Trompeuse sur déséquilibre 75/25"],
     ["Précision",   "—",     "0.249", "0.245", "0.246", "≈ 25% → prédit bien la minorité ?"],
     ["Rappel",      "0.000", "0.509", "0.355", "0.441", "LR détecte ~50% des cas à risque"],
     ["F1",          "0.000", "0.335", "0.290", "0.316", "Métrique équilibrée → faible"],
     ["AUC-ROC",     "0.500", "0.499", "0.494", "0.499", "≈ 0.5 = tirage aléatoire"]],
    fs=12, cw=[Inches(1.25), Inches(1.35), Inches(1.35), Inches(1.7), Inches(1.35), Inches(5.6)])

block(sl, Inches(0.35), CONTENT_Y + Inches(4.55), W - Inches(0.7), Inches(0.95),
    "Interprétation des résultats",
    ["Modèle naïf 'tout à 0' : accuracy = 75.2% sans apprendre → accuracy seule trompeuse  (Sokolova & Lapalme, 2009)",
     "AUC ≈ 0.5 pour tous = aucun pouvoir discriminant — équivalent tirage aléatoire",
     "Cause : disease_risk généré sans dépendance aux autres variables dans le dataset synthétique"])

img(sl, FIGS+"roc_curves.png", Inches(0.35), CONTENT_Y + Inches(5.7), w=Inches(6.0))
img(sl, FIGS+"confusion_matrices.png", Inches(6.6), CONTENT_Y + Inches(5.7), w=Inches(6.4))

# ── D2 : SMOTE & class_weight ─────────────────────────────────
sl = new_slide()
title_bar(sl, "D — Gestion du déséquilibre : class_weight & SMOTE", 16)

txt(sl, Inches(0.35), CONTENT_Y + Inches(0.1), W - Inches(0.7), Inches(0.28),
    "Déséquilibre : 75.2% sans risque (classe 0)  vs  24.8% à risque (classe 1)  →  deux stratégies comparées",
    size=14, bold=True, color=BLEU)

# Class weight
rect(sl, Inches(0.35), CONTENT_Y + Inches(0.55), Inches(5.9), Inches(0.28), BLEU)
txt(sl, Inches(0.5), CONTENT_Y + Inches(0.58), Inches(5.6), Inches(0.22),
    "Stratégie 1 — class_weight='balanced'", size=14, bold=True, color=BLANC)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(0.97), Inches(5.9), Inches(0.88),
    "wₖ = n / (K × nₖ)",
    "n = total observations  |  K = nb de classes  |  nₖ = taille de la classe k\n"
    "→ classe minoritaire (risk=1) reçoit un poids 3x plus élevé pendant l'entraînement")

bullets(sl, Inches(0.35), CONTENT_Y + Inches(2.0), Inches(5.9), Inches(1.7),
    ["Modifie la FONCTION DE PERTE, pas les données",
     "Pénalise davantage les erreurs sur la classe minoritaire",
     "Pas de surapprentissage additionnel",
     "Rappel monte mais précision baisse → trade-off"],
    b_size=13)

# SMOTE
rect(sl, Inches(6.5), CONTENT_Y + Inches(0.55), Inches(6.45), Inches(0.28), BLEU)
txt(sl, Inches(6.65), CONTENT_Y + Inches(0.58), Inches(6.1), Inches(0.22),
    "Stratégie 2 — SMOTE  (Chawla et al., 2002)", size=14, bold=True, color=BLANC)

formula_box(sl, Inches(6.5), CONTENT_Y + Inches(0.97), Inches(6.45), Inches(0.88),
    "x_new = xᵢ + λ × (x_nn − xᵢ)   avec λ ∈ [0, 1]",
    "xᵢ : exemple minoritaire  |  x_nn : k-plus proche voisin dans la même classe\n"
    "→ création synthétique de nouveaux exemples par interpolation linéaire")

bullets(sl, Inches(6.5), CONTENT_Y + Inches(2.0), Inches(6.45), Inches(1.7),
    ["Modifie les DONNÉES (sur-échantillonnage)",
     "Crée de vrais nouveaux exemples — pas de simples doublons",
     "Appliqué sur le jeu d'entraînement SEULEMENT",
     "Peut générer du bruit si classes mal séparées"],
    b_size=13)

add_table(sl, Inches(0.35), CONTENT_Y + Inches(3.9),
          W - Inches(0.7), Inches(1.05),
    ["Stratégie", "Rappel", "Précision", "F1", "AUC", "Conclusion"],
    [["Sans correction",      "0.355", "0.245", "0.290", "0.494", "Baseline"],
     ["class_weight balanced","0.680", "0.205", "0.315", "0.492", "Rappel↑ mais précision↓↓"],
     ["SMOTE",                "0.420", "0.242", "0.307", "0.498", "Résultat similaire à baseline"]],
    fs=13, cw=[Inches(2.5), Inches(1.1), Inches(1.2), Inches(0.9), Inches(0.9), Inches(6.3)])

block(sl, Inches(0.35), CONTENT_Y + Inches(5.15), W - Inches(0.7), Inches(1.0),
    "Conclusion — déséquilibre vs absence de signal",
    ["Les deux stratégies confirment : AUC ≈ 0.5 indépendamment de la gestion du déséquilibre",
     "Le problème N'EST PAS le déséquilibre des classes — c'est l'absence totale de signal prédictif",
     "Ni class_weight ni SMOTE ne peuvent créer un signal qui n'existe pas dans les données"])


# ══════════════════════════════════════════════════════════════
# SECTION E — CLUSTERING
# ══════════════════════════════════════════════════════════════

sl = new_slide()
section_header(sl, 17, "E", "Clustering")

# ── E1 : K-Means ──────────────────────────────────────────────
sl = new_slide()
title_bar(sl, "E — K-Means : algorithme & formules", 18)

bullets(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(5.9), Inches(1.5),
    ["Algorithme de partitionnement (MacQueen, 1967)",
     "Minimise l'inertie intra-cluster (dispersion totale)",
     "Deux étapes alternées jusqu'à convergence",
     "100 000 individus — O(n) rapide"],
    title="K-Means — Principe", b_size=14)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(1.75), Inches(5.9), Inches(1.0),
    "Distance euclidienne : d(x, y) = √[ Σⱼ (xⱼ − yⱼ)² ]",
    "p = nombre de variables (ici p=15). Sans standardisation, daily_steps\n"
    "dominerait sleep_hours → normalisation z-score obligatoire")

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(2.9), Inches(5.9), Inches(1.0),
    "Affectation : cluster(xᵢ) = argmin_j  ||xᵢ − μⱼ||²",
    "Chaque individu est affecté au centroïde le plus proche")

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(4.05), Inches(5.9), Inches(1.0),
    "Mise à jour : μⱼ = (1/|Cⱼ|) × Σ_{x∈Cⱼ} x",
    "Centroïde = vecteur des moyennes de toutes les variables dans le cluster")

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(5.2), Inches(5.9), Inches(1.05),
    "Objectif : W(k) = Σⱼ Σ_{x∈Cⱼ} ||x − μⱼ||²   (inertie intra-cluster)",
    "Converge vers un minimum local — 10 initialisations aléatoires (n_init=10)")

add_table(sl, Inches(6.5), CONTENT_Y + Inches(0.1),
          Inches(6.45), Inches(1.7),
    ["Cluster", "Taille", "IMC moyen", "Systolique", "Profil"],
    [["0", "32 754", "35.0 (obèse)",    "142 mmHg (HTA)", "Obèse hypertendu"],
     ["1", "27 673", "29.1 (surpoids)", "115 mmHg (OK)",  "Surpoids normotendu"],
     ["2", "39 573", "24.0 (normal)",   "142 mmHg (HTA)", "Normal hypertendu"]],
    fs=12, cw=[Inches(0.9), Inches(1.1), Inches(1.55), Inches(1.55), Inches(1.3)])

txt(sl, Inches(6.5), CONTENT_Y + Inches(1.95), Inches(6.45), Inches(0.28),
    "Partition dominée par IMC (σ=6.3) et systolic_bp (σ=26.0) → variables à plus grande variance",
    size=12, italic=True, color=RGBColor(80,80,80))

img(sl, FIGS+"kmeans_selection_k.png", Inches(6.5), CONTENT_Y + Inches(2.35), w=Inches(3.0))
img(sl, FIGS+"kmeans_pca.png", Inches(9.65), CONTENT_Y + Inches(2.35), w=Inches(3.3))
cap(sl, Inches(6.5), H - Inches(0.52), Inches(3.0),
    "Inertie & silhouette vs k  →  k=3 optimal")
cap(sl, Inches(9.65), H - Inches(0.52), Inches(3.3),
    "Clusters K-Means — ACP 2D (24.3%)")

# ── E2 : Sélection k + Silhouette ─────────────────────────────
sl = new_slide()
title_bar(sl, "E — Sélection de k : inertie & coefficient de silhouette", 19)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(6.1), Inches(0.88),
    "Inertie : W(k) = Σ_{j=1}^{k} Σ_{x∈Cⱼ} ||x − μⱼ||²",
    "Décroît quand k augmente — on cherche le 'coude' où le gain devient marginal")

formula_box(sl, Inches(6.55), CONTENT_Y + Inches(0.1), Inches(6.45), Inches(0.88),
    "Silhouette : s(i) = (b(i) − a(i)) / max(a(i), b(i))",
    "a(i) = distance moy. aux membres de son cluster  |  b(i) = distance moy. au cluster voisin le + proche")

add_table(sl, Inches(0.35), CONTENT_Y + Inches(1.15),
          Inches(5.7), Inches(2.4),
    ["k", "Inertie W(k)", "Silhouette s", "Interprétation"],
    [["2",  "1 346 768", "0.0979", "Acceptable"],
     ["3",  "1 238 076", "0.1030", "★ Meilleur — retenu"],
     ["4",  "1 188 303", "0.0967", "Silhouette baisse"],
     ["5",  "1 152 260", "0.0886", "Gain inertie faible"],
     ["10", "1 031 548", "0.0744", "Déclin clair"]],
    fs=13, cw=[Inches(0.55), Inches(1.8), Inches(1.4), Inches(1.9)])

txt(sl, Inches(0.35), CONTENT_Y + Inches(3.7), Inches(5.7), Inches(0.5),
    "k=3 maximise la silhouette (0.1030) ET correspond au coude de l'inertie  →  k=3 retenu",
    size=13, bold=True, color=BLEU)

block(sl, Inches(0.35), CONTENT_Y + Inches(4.35), Inches(5.7), Inches(2.0),
    "Interprétation de la silhouette",
    ["s ≈ 1 : individu bien assigné — loin des autres clusters",
     "s ≈ 0 : individu en frontière entre deux clusters",
     "s < 0 : individu probablement mal assigné",
     "Seuils : >0.7 excellent · 0.5–0.7 bon · <0.5 faible",
     "Ici 0.103 : faible — attendu sur données sans structure causale"])

formula_box(sl, Inches(6.55), CONTENT_Y + Inches(1.15), Inches(6.45), Inches(1.1),
    "Davies-Bouldin = (1/k) × Σⱼ max_{i≠j} [(σᵢ + σⱼ) / d(μᵢ, μⱼ)]",
    "σⱼ = dispersion intra-cluster j  |  d(μᵢ,μⱼ) = distance entre centroïdes\n"
    "Plus petit est meilleur (vise 0). Ici DB ≈ 2.66 pour K-Means.")

add_table(sl, Inches(6.55), CONTENT_Y + Inches(2.45),
          Inches(6.45), Inches(1.25),
    ["Méthode", "Silhouette", "Davies-Bouldin", "n individus"],
    [["K-Means (k=3)",  "0.103", "2.655", "100 000"],
     ["CAH Ward (k=3)", "0.105", "2.721", "5 000"]],
    fs=13, cw=[Inches(2.1), Inches(1.4), Inches(1.7), Inches(1.2)])

txt(sl, Inches(6.55), CONTENT_Y + Inches(3.85), Inches(6.45), Inches(0.35),
    "CAH légèrement meilleure en silhouette (0.105 > 0.103) mais sous-échantillon 5000 individus → K-Means privilégié sur données complètes",
    size=12, italic=True, color=RGBColor(80,80,80))

img(sl, FIGS+"kmeans_pca.png", Inches(6.55), CONTENT_Y + Inches(4.4), w=Inches(6.45))
cap(sl, Inches(6.55), H - Inches(0.4), Inches(6.45),
    "Projection ACP — variance expliquée : 24.3% (2 composantes sur 15 dimensions quasi-indépendantes)")

# ── E3 : CAH Ward ─────────────────────────────────────────────
sl = new_slide()
title_bar(sl, "E — CAH : Classification Hiérarchique Ascendante (liaison Ward)", 20)

bullets(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(5.9), Inches(1.5),
    ["Principe : chaque individu commence dans son propre cluster",
     "À chaque étape : fusion des 2 clusters dont la fusion est la moins coûteuse",
     "Produit un dendrogramme — arbre de fusion lisible visuellement",
     "Complexité O(n²) → appliqué sur 5 000 individus (sous-échantillon)"],
    title="CAH — Principe", b_size=14)

formula_box(sl, Inches(0.35), CONTENT_Y + Inches(1.75), Inches(5.9), Inches(1.2),
    "Coût Ward : Δ(A,B) = (|A|×|B|)/(|A|+|B|) × ||μ_A − μ_B||²",
    "|A|, |B| = tailles des clusters  |  μ_A, μ_B = centroïdes\n"
    "Minimise l'augmentation de variance intra-cluster à chaque fusion.\n"
    "→ clusters compacts et de tailles équilibrées")

bullets(sl, Inches(0.35), CONTENT_Y + Inches(3.15), Inches(5.9), Inches(1.5),
    ["Single linkage : distance entre points les plus proches → clusters en chaîne",
     "Complete linkage : distance entre points les plus éloignés → clusters sphériques",
     "Average linkage : distance moyenne → compromis",
     "Ward : minimise la variance → le plus adapté aux données continues"],
    title="Comparaison des méthodes de liaison", b_size=13)

add_table(sl, Inches(0.35), CONTENT_Y + Inches(4.85),
          Inches(5.9), Inches(1.6),
    ["Cluster", "Taille", "% Fumeurs", "Systolique", "Profil CAH"],
    [["0", "2 928", "2%",    "142 mmHg", "Non-fumeurs hypertendus"],
     ["1", "1 407", "18%",   "115 mmHg", "Fumeurs modérés, normotendus"],
     ["2", "665",  "100%",  "142 mmHg", "Fumeurs intenses, hypertendus"]],
    fs=12, cw=[Inches(0.8), Inches(0.9), Inches(1.1), Inches(1.35), Inches(1.7)])

txt(sl, Inches(0.35), CONTENT_Y + Inches(6.6), Inches(5.9), Inches(0.28),
    "CAH isole les fumeurs intenses (100%) — groupe dilué par K-Means → méthodes complémentaires",
    size=12, italic=True, color=RGBColor(80,80,80))

block(sl, Inches(6.5), CONTENT_Y + Inches(0.1), Inches(6.45), Inches(2.6),
    "SHAP — formule complète  (Lundberg & Lee, 2017)",
    ["Valeur de Shapley pour la variable j :",
     "φⱼ = Σ_{S⊆F\\{j}} [|S|!(|F|−|S|−1)! / |F|!] × [f(S∪{j}) − f(S)]",
     "F = ensemble des variables  |  S = sous-ensemble sans j",
     "f(S) = prédiction du modèle avec les variables de S seulement",
     "Somme sur toutes les coalitions possibles de variables",
     "Propriété : φⱼ > 0 → variable pousse vers le risque ; φⱼ < 0 → vers le non-risque"])

img(sl, FIGS+"shap_summary.png", Inches(6.5), CONTENT_Y + Inches(2.8), w=Inches(6.45))
cap(sl, Inches(6.5), H - Inches(0.45), Inches(6.45),
    "SHAP summary plot — XGBoost  (divergence avec feature importance Gini = résultat analytique sur données sans signal)")


# ══════════════════════════════════════════════════════════════
# Sauvegarde
# ══════════════════════════════════════════════════════════════
OUT = ("/Users/ardachammahamatteguene/Desktop/"
       "Bureau - MacBook Pro de Ardacham/Memoire/presentation_theorie.pptx")
prs.save(OUT)
print(f"Saved: {OUT}  — {prs.slides.__len__()} slides")
