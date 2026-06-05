"""
Génère presentation.pptx — design fidèle au Beamer latex/presentation.pdf
  • Fond blanc
  • Titre de diapo : bande grisclair (240,245,250) + texte bleu foncé (0,70,140)
  • Trait bleu accent (0,102,204) sous le titre  — exactement comme le \hrule Beamer
  • Blocs : en-tête bleu foncé / texte blanc  +  corps grisclair / texte noir
  • Puces sobres, texte noir sur fond blanc
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette identique au Beamer ───────────────────────────────
BLEU       = RGBColor(0,   70, 140)   # bleutitre  — titres
BLEU_ACC   = RGBColor(0,  102, 204)   # bleuaccent — trait + sous-titres
GRIS_CLAIR = RGBColor(240, 245, 250)  # grisclair  — fond titre + corps bloc
NOIR       = RGBColor(0,   0,   0)
BLANC      = RGBColor(255, 255, 255)
ROUGE_AV   = RGBColor(180,  30,  30)  # alerte

# ── Dimensions 16:9 ──────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)
FIGS = ("/Users/ardachammahamatteguene/Desktop/"
        "Bureau - MacBook Pro de Ardacham/Memoire/latex/figures/")

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────
# Primitives
# ─────────────────────────────────────────────────────────────

def new_slide():
    sl = prs.slides.add_slide(BLANK)
    # fond blanc
    bg = sl.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLANC
    return sl

def add_rect(sl, x, y, w, h, fill, border=None):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s

def add_text(sl, x, y, w, h, text, size=14, bold=False, italic=False,
             color=NOIR, align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text         = text
    r.font.size    = Pt(size)
    r.font.bold    = bold
    r.font.italic  = italic
    r.font.color.rgb = color
    return tb

def add_img(sl, path, x, y, w=None, h=None):
    if w and h: sl.shapes.add_picture(path, x, y, w, h)
    elif w:     sl.shapes.add_picture(path, x, y, width=w)
    elif h:     sl.shapes.add_picture(path, x, y, height=h)
    else:       sl.shapes.add_picture(path, x, y)

# ─────────────────────────────────────────────────────────────
# Composants haut niveau
# ─────────────────────────────────────────────────────────────

TITLE_H    = Inches(0.82)
HRULE_H    = Pt(2)
CONTENT_Y  = TITLE_H + Inches(0.12)

def title_bar(sl, text):
    """
    Bande grisclair + texte bleu foncé bold + trait bleuaccent en dessous
    — reproduction exacte du frametitle Beamer.
    """
    add_rect(sl, 0, 0, W, TITLE_H, GRIS_CLAIR)
    add_text(sl, Inches(0.35), Inches(0.14), W - Inches(0.7), Inches(0.58),
             text, size=24, bold=True, color=BLEU)
    add_rect(sl, 0, TITLE_H, W, HRULE_H, BLEU_ACC)

def bullet_list(sl, x, y, w, h, items, title=None, size=14, t_size=15,
                t_color=BLEU, b_color=NOIR):
    """Liste à puces sobre (tiret –  noir ou bleu)."""
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]; first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size  = Pt(t_size)
        r.font.bold  = True
        r.font.color.rgb = t_color
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "–  " + item
        r.font.size  = Pt(size)
        r.font.color.rgb = b_color

def beamer_block(sl, x, y, w, h, title, items, t_size=14, b_size=13, t_col=BLEU):
    """
    Bloc style Beamer : en-tête couleur pleine / texte blanc
                        corps grisclair / texte noir.
    """
    HDR = Inches(0.38)
    add_rect(sl, x, y, w, HDR, t_col)
    add_text(sl, x + Inches(0.1), y + Inches(0.04),
             w - Inches(0.2), HDR - Inches(0.06),
             title, size=t_size, bold=True, color=BLANC)
    add_rect(sl, x, y + HDR, w, h - HDR, GRIS_CLAIR)
    tb = sl.shapes.add_textbox(x + Inches(0.12), y + HDR + Inches(0.05),
                                w - Inches(0.24), h - HDR - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = "–  " + item
        r.font.size = Pt(b_size)
        r.font.color.rgb = NOIR

def add_table(sl, x, y, w, h, headers, rows, fs=12, cw=None):
    """Tableau : en-tête BLEU/blanc, lignes alternées blanc/grisclair."""
    nc = len(headers)
    t  = sl.shapes.add_table(len(rows) + 1, nc, x, y, w, h).table
    if cw:
        for i, c in enumerate(cw): t.columns[i].width = c
    for j, hdr in enumerate(headers):
        cell = t.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = BLEU
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hdr; r.font.size = Pt(fs); r.font.bold = True
        r.font.color.rgb = BLANC
    for i, row in enumerate(rows):
        bg = BLANC if i % 2 == 0 else GRIS_CLAIR
        for j, val in enumerate(row):
            cell = t.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val); r.font.size = Pt(fs)
            r.font.color.rgb = NOIR
    return t

def caption(sl, x, y, w, text):
    add_text(sl, x, y, w, Inches(0.32), text, size=11,
             italic=True, color=RGBColor(100, 100, 100), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# DIAPO 1 — TITRE
# ══════════════════════════════════════════════════════════════
sl = new_slide()

# Ligne déco bleu accent (comme le trait Beamer)
add_rect(sl, 0, Inches(2.2), W, Pt(2), BLEU_ACC)
add_rect(sl, 0, Inches(5.1), W, Pt(2), BLEU_ACC)

add_text(sl, Inches(0.8), Inches(2.4), W - Inches(1.6), Inches(1.3),
         "Modélisation des comportements de santé\npar apprentissage automatique",
         size=34, bold=True, color=BLEU, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(3.78), W - Inches(1.6), Inches(0.55),
         "Régression, classification et clustering sur données synthétiques",
         size=18, italic=True, color=BLEU_ACC, align=PP_ALIGN.CENTER)

add_text(sl, Inches(0.8), Inches(5.25), W - Inches(1.6), Inches(0.4),
         "Ardacham Mahamat Teguene",
         size=16, bold=True, color=NOIR, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(5.68), W - Inches(1.6), Inches(0.35),
         "Master 1 MIASHS  —  Université de Montpellier Paul Valéry",
         size=14, color=RGBColor(80, 80, 80), align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(6.05), W - Inches(1.6), Inches(0.35),
         "Tuteurs : Sophie Lèbre & Jérôme Pasquet  —  Juin 2026",
         size=13, color=RGBColor(120, 120, 120), align=PP_ALIGN.CENTER)

# Numéro de diapo style Beamer (bas droite)
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "1", size=10, color=RGBColor(140, 140, 140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 2 — PLAN
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Plan de la présentation")

sections = [
    "Contexte et problématique",
    "Jeu de données & pipeline CRISP-DM",
    "Analyse exploratoire (EDA)",
    "Modélisation supervisée — Régression",
    "Modélisation supervisée — Classification",
    "Clustering et interprétabilité",
    "Bilan & Perspectives",
]
for i, sec in enumerate(sections):
    y = CONTENT_Y + Inches(0.18) + i * Inches(0.74)
    add_rect(sl, Inches(0.35), y, Inches(0.42), Inches(0.52), BLEU)
    add_text(sl, Inches(0.35), y + Inches(0.07), Inches(0.42), Inches(0.38),
             str(i+1), size=17, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(0.95), y + Inches(0.09), Inches(11), Inches(0.38),
             sec, size=17, color=NOIR)

add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "2", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 3 — CONTEXTE & PROBLÉMATIQUE
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Contexte et problématique")

# Bloc problématique (fond grisclair + bordure gauche bleuaccent)
add_rect(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
         W - Inches(0.7), Inches(0.95), GRIS_CLAIR)
add_rect(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
         Pt(4), Inches(0.95), BLEU_ACC)
add_text(sl, Inches(0.6), CONTENT_Y + Inches(0.15),
         W - Inches(1.1), Inches(0.25),
         "PROBLÉMATIQUE", size=12, bold=True, color=BLEU_ACC)
add_text(sl, Inches(0.6), CONTENT_Y + Inches(0.4),
         W - Inches(1.1), Inches(0.58),
         "Quels facteurs du mode de vie influencent le plus l'état de santé "
         "— et peut-on les identifier, les modéliser et les expliquer de manière rigoureuse ?",
         size=14, bold=True, color=BLEU)

y2 = CONTENT_Y + Inches(1.2)
bullet_list(sl, Inches(0.35), y2, Inches(5.9), Inches(2.1),
    ["Les maladies chroniques résultent de comportements quotidiens mesurables",
     "ML peut capturer des relations complexes inaccessibles à la stat. classique",
     "Double enjeu : prédire ET expliquer (interprétabilité)"],
    title="Contexte data science & santé", size=14)

bullet_list(sl, Inches(6.7), y2, Inches(6.2), Inches(2.1),
    ["Obésité : IMC > 30  (seuil OMS)",
     "Hypercholestérolémie : cholestérol > 200 mg/dL  (NCEP)",
     "Hypertension : pression systolique ≥ 140 mmHg  (ESC)"],
    title="Seuils médicaux de référence", size=14)

add_text(sl, Inches(0.35), H - Inches(0.9), W - Inches(0.7), Inches(0.3),
         "Variable cible régression : cholesterol (mg/dL)  |  "
         "Variable cible classification : disease_risk (0/1)",
         size=13, italic=True, color=RGBColor(80,80,80))
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "3", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 4 — JEU DE DONNÉES
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Jeu de données — Health & Lifestyle Dataset")

add_table(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
          Inches(5.7), Inches(3.05),
    ["Type", "Variables"],
    [["Continues (8)", "age, bmi, daily_steps, sleep_hours, cholesterol,\nsystolic_bp, diastolic_bp, resting_hr"],
     ["Binaires (4)", "smoker, alcohol, family_history, disease_risk"],
     ["Catégorielle (1)", "gender"],
     ["Cible régression", "cholesterol (mg/dL)"],
     ["Cible classification", "disease_risk (0 = sain / 1 = à risque)"]],
    fs=13, cw=[Inches(1.85), Inches(3.85)])

bullet_list(sl, Inches(0.35), CONTENT_Y + Inches(3.3), Inches(5.7), Inches(2.0),
    ["100 000 individus, 15 variables",
     "Dataset synthétique (Kaggle) — sans contrainte RGPD",
     "Aucune valeur manquante, aucun outlier (IQR)",
     "Distributions quasi-uniformes (skewness < |0.02|)"],
    title="Points clés", size=14)

add_img(sl, FIGS+"disease_risk_distribution.png",
        Inches(6.4), CONTENT_Y + Inches(0.05), w=Inches(6.55))
caption(sl, Inches(6.4), H - Inches(0.72), Inches(6.55),
        "Déséquilibre 75% / 25%  — un modèle naïf atteint 75% d'accuracy sans rien apprendre")
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "4", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 5 — PIPELINE CRISP-DM
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Pipeline CRISP-DM")

etapes = [
    ("EDA",          "Corrélations < 0.01\nDéséquilibre 75/25\nPas d'outliers"),
    ("Feature\nEng.","Encodage gender\n3 variables métier\nLeakage détecté"),
    ("Régression",   "6 modèles\nRandSearchCV\nK-Fold k=5"),
    ("Classif.",     "3 modèles\nSMOTE vs CW\nStrat. K-Fold"),
    ("Clustering",   "K-Means k=3\nCAH Ward\nACP 24.3%"),
    ("Interprét.",   "Coeff. LR/Ridge\nFI Gini RF/XGB\nSHAP"),
]
bw = Inches(1.88); bh = Inches(1.45); gap = Inches(0.2)
sx = Inches(0.35)
by = CONTENT_Y + Inches(0.15)

for i, (titre, desc) in enumerate(etapes):
    bx = sx + i * (bw + gap)
    add_rect(sl, bx, by, bw, Inches(0.4), BLEU)
    add_text(sl, bx, by + Inches(0.04), bw, Inches(0.33),
             titre, size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_rect(sl, bx, by + Inches(0.4), bw, bh - Inches(0.4), GRIS_CLAIR)
    add_text(sl, bx + Inches(0.08), by + Inches(0.46),
             bw - Inches(0.16), bh - Inches(0.52),
             desc, size=11, color=NOIR)
    if i < len(etapes) - 1:
        ax = bx + bw + Inches(0.04)
        ay = by + bh / 2
        add_rect(sl, ax, ay - Pt(2), gap - Inches(0.06), Pt(4), BLEU_ACC)

y3 = CONTENT_Y + Inches(1.8)
add_text(sl, Inches(0.35), y3, W - Inches(0.7), Inches(0.3),
         "Approche itérative : chaque étape informe la suivante "
         "(EDA → anticipation des problèmes de signal → choix des modèles → interprétabilité)",
         size=13, italic=True, color=RGBColor(80,80,80))

add_table(sl, Inches(0.35), y3 + Inches(0.4), W - Inches(0.7), Inches(2.5),
    ["Tâche", "Meilleur modèle", "Métrique principale", "Conclusion"],
    [["Régression",    "XGBoost tuné",       "R² = 0.0001",     "Aucun signal — tous équivalents"],
     ["Classification","Logistic Regression", "AUC = 0.499",     "Équivalent tirage aléatoire"],
     ["Clustering",    "K-Means k=3",         "Silhouette = 0.103","Profils biologiques interprétables"]],
    fs=13, cw=[Inches(2.0), Inches(2.5), Inches(2.5), Inches(5.8)])
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "5", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 6 — EDA
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Analyse exploratoire (EDA)")

add_text(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
         Inches(7.9), Inches(0.28),
         "Distributions des variables continues", size=14, bold=True, color=BLEU)
add_img(sl, FIGS+"distributions_continues.png",
        Inches(0.35), CONTENT_Y + Inches(0.42), w=Inches(7.9))

add_text(sl, Inches(8.5), CONTENT_Y + Inches(0.1),
         Inches(4.5), Inches(0.28),
         "Corrélations avec la cible", size=14, bold=True, color=BLEU)
add_img(sl, FIGS+"correlations_target.png",
        Inches(8.5), CONTENT_Y + Inches(0.42), w=Inches(4.5))

add_rect(sl, Inches(0.35), H - Inches(1.4), W - Inches(0.7), Inches(1.05), GRIS_CLAIR)
add_rect(sl, Inches(0.35), H - Inches(1.4), Pt(4), Inches(1.05), ROUGE_AV)
add_text(sl, Inches(0.6), H - Inches(1.35), W - Inches(1.0), Inches(0.25),
         "Résultat central", size=12, bold=True, color=ROUGE_AV)
add_text(sl, Inches(0.6), H - Inches(1.1), W - Inches(1.0), Inches(0.68),
         "Toutes les corrélations < |0.01|. La plus élevée avec disease_risk : "
         "resting_hr avec r = 0.005 (non significatif sur 100 000 individus). "
         "Conclusion : disease_risk généré indépendamment → justifie l'exploration de méthodes non-linéaires.",
         size=13, color=NOIR)
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "6", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 7 — FEATURE ENGINEERING
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Feature Engineering — 16 variables explicatives")

bullet_list(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(5.9), Inches(1.35),
    ["gender → gender_enc (0/1)",
     "Suppression de id (pas de pouvoir prédictif)"],
    title="Encodage des variables catégorielles", size=14)

bullet_list(sl, Inches(0.35), CONTENT_Y + Inches(1.6), Inches(5.9), Inches(2.4),
    ["hypertension : systolic_bp ≥ 140 ou diastolic_bp ≥ 90  (seuil ESC)",
     "bmi_cat : 0=sous-poids / 1=normal / 2=surpoids / 3=obèse  (seuil OMS)",
     "lifestyle_score : score composite  (steps + sleep + water  – smoker – alcohol – bmi – calories)"],
    title="3 variables métier créées à partir des connaissances médicales", size=14)

bullet_list(sl, Inches(0.35), CONTENT_Y + Inches(4.15), Inches(5.9), Inches(1.35),
    ["Standardisation z-score pour les modèles linéaires",
     "Normalisation min-max pour K-Means (distance euclidienne)"],
    title="Mise à l'échelle", size=14)

beamer_block(sl, Inches(6.5), CONTENT_Y + Inches(0.1),
             Inches(6.45), Inches(2.6),
    "⚠  Data Leakage détecté et corrigé",
    ["hypercholesterol = 1 si cholesterol > 200 mg/dL",
     "Directement dérivée de la variable cible → 'donner la réponse' au modèle",
     "R² artificiel = 0.675 avec leakage → R² ≈ 0 après suppression",
     "Exclue de la régression, conservée uniquement pour le clustering"],
    t_col=ROUGE_AV)

add_text(sl, Inches(6.5), CONTENT_Y + Inches(2.85), Inches(6.45), Inches(0.3),
         "Dataset final : 15 variables originales → 16 features explicatives  (14 − 1 + 3 = 16)",
         size=14, bold=True, color=BLEU)
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "7", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 8 — RÉGRESSION : PROTOCOLE & MODÈLES
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Régression — Protocole & justification des modèles")

bullet_list(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(5.9), Inches(1.6),
    ["Split 80/20  (random_state=42) — 80 000 train / 20 000 test",
     "K-Fold k=5 : 16 000 exemples/fold, bon compromis coût/fiabilité",
     "RandomizedSearchCV n_iter=20 : 20 tirages aléatoires vs GridSearch exhaustif",
     "Métriques : R², RMSE (mg/dL), MAE (mg/dL)"],
    title="Protocole expérimental", size=14)

add_table(sl, Inches(0.35), CONTENT_Y + Inches(1.85),
          Inches(5.9), Inches(3.6),
    ["Modèle", "Justification"],
    [["Régression linéaire (baseline)", "Référence obligatoire — plancher de performance"],
     ["Ridge  (L2)",    "Réduit les coeff. sans les annuler — stabilise si variables peu informatives"],
     ["Lasso  (L1)",    "Force des coeff. à 0 → sélection automatique de variables"],
     ["ElasticNet",     "Combine L1 + L2 — robuste quand on ne sait pas lequel convient"],
     ["Random Forest",  "Ensembliste non-linéaire parallèle — capte les interactions"],
     ["XGBoost",        "Ensembliste séquentiel — corrige les erreurs itérativement"]],
    fs=12, cw=[Inches(2.15), Inches(3.75)])

beamer_block(sl, Inches(6.5), CONTENT_Y + Inches(0.1),
             Inches(6.45), Inches(2.15),
    "Logique : spectre de complexité croissante",
    ["Baseline linéaire → plancher de performance",
     "Régularisées → teste si contraindre les coeff. aide",
     "Ensemblistes → capture les non-linéarités potentielles",
     "Si complexe ≤ baseline → complexité non justifiée (rasoir d'Occam)"])

add_img(sl, FIGS+"distributions_continues.png",
        Inches(6.5), CONTENT_Y + Inches(2.4), w=Inches(6.45))
caption(sl, Inches(6.5), H - Inches(0.52), Inches(6.45),
        "Distributions quasi-uniformes — aucune transformation nécessaire")
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "8", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 9 — RÉGRESSION : RÉSULTATS
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Régression — Résultats")

add_table(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
          Inches(7.6), Inches(2.8),
    ["Modèle", "R²", "RMSE", "MAE", "CV-R²"],
    [["Linéaire (baseline)",  "0.0001",   "43.33", "37.59", "≈ 0"],
     ["Ridge (λ=7197, tuné)", "0.0001",   "43.33", "37.59", "≈ 0"],
     ["Lasso (λ=9.54, tuné)", "−0.000",   "43.33", "37.60", "≈ 0"],
     ["ElasticNet",           "−0.000",   "43.33", "37.60", "≈ 0"],
     ["Random Forest (tuné)", "−0.005",   "43.44", "37.67", "≈ 0"],
     ["XGBoost (tuné)",       "0.0001",   "43.33", "37.59", "≈ 0"]],
    fs=13, cw=[Inches(2.5), Inches(1.0), Inches(1.35), Inches(1.35), Inches(1.3)])

beamer_block(sl, Inches(0.35), CONTENT_Y + Inches(3.05),
             Inches(7.6), Inches(2.85),
    "Analyse critique",
    ["Lasso (tuné) : 16/16 coefficients = 0 → prédit uniquement la moyenne (224.3 mg/dL)",
     "Random Forest non-tuné : R² < 0 → pire que la moyenne (surajustement du bruit)",
     "Data leakage corrigé : R² 0.675 → ≈ 0 après suppression de hypercholesterol",
     "Différences entre −0.005 et 0.0001 statistiquement négligeables → aucun modèle retenu",
     "Cause : variables générées indépendamment — signal absent par construction"])

add_img(sl, FIGS+"correlation_matrix.png",
        Inches(8.1), CONTENT_Y + Inches(0.1), w=Inches(4.9))
caption(sl, Inches(8.1), H - Inches(0.52), Inches(4.9),
        "Matrice de corrélation — toutes valeurs < |0.01|")
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "9", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 10 — CLASSIFICATION : RÉSULTATS
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Classification — Résultats")

add_table(sl, Inches(0.35), CONTENT_Y + Inches(0.1),
          Inches(6.55), Inches(2.15),
    ["Modèle", "Accuracy", "Précision", "Rappel", "F1", "AUC"],
    [["Naïf (tout à 0)",     "0.752", "n/a",   "0.000", "0.000", "0.500"],
     ["Logistic Regression", "0.498", "0.249", "0.509", "0.335", "0.499"],
     ["Random Forest",       "0.568", "0.245", "0.355", "0.290", "0.494"],
     ["XGBoost",             "0.525", "0.246", "0.441", "0.316", "0.499"]],
    fs=13, cw=[Inches(2.2), Inches(0.82), Inches(0.82), Inches(0.82), Inches(0.72), Inches(0.8)])

beamer_block(sl, Inches(0.35), CONTENT_Y + Inches(2.4),
             Inches(6.55), Inches(2.85),
    "Analyse critique",
    ["AUC ≈ 0.5 pour tous = équivalent à un tirage aléatoire",
     "class_weight='balanced' : rappel monte mais accuracy < 50% → sur-correction",
     "SMOTE : même résultat — ne peut pas créer un signal absent",
     "Accuracy seule trompeuse sur données déséquilibrées → F1 + AUC obligatoires",
     "Cohérent avec régression : disease_risk généré indépendamment de toutes les variables"])

add_img(sl, FIGS+"roc_curves.png", Inches(7.1), CONTENT_Y + Inches(0.1), w=Inches(3.2))
caption(sl, Inches(7.1), CONTENT_Y + Inches(3.5), Inches(3.2),
        "Courbes ROC — AUC ≈ 0.5")
add_img(sl, FIGS+"confusion_matrices.png", Inches(10.3), CONTENT_Y + Inches(0.1), w=Inches(2.7))
caption(sl, Inches(10.3), CONTENT_Y + Inches(3.5), Inches(2.7),
        "Matrices de confusion")
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "10", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 11 — CLUSTERING
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Clustering — K-Means & CAH")

add_text(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Inches(5.6), Inches(0.28),
         "K-Means (k=3, 100 000 individus)", size=14, bold=True, color=BLEU)
add_table(sl, Inches(0.35), CONTENT_Y + Inches(0.43),
          Inches(5.6), Inches(1.5),
    ["Cluster", "Profil", "IMC moy.", "Systolique"],
    [["0", "Obèse hypertendu",    "35.0", "142 mmHg"],
     ["1", "Surpoids normotendu", "29.1", "115 mmHg"],
     ["2", "Normal hypertendu",   "24.0", "142 mmHg"]],
    fs=13, cw=[Inches(0.78), Inches(2.17), Inches(1.25), Inches(1.35)])
add_text(sl, Inches(0.35), CONTENT_Y + Inches(2.05), Inches(5.6), Inches(0.3),
         "Partition par IMC + pression systolique (variables à plus grande variance)",
         size=12, italic=True, color=RGBColor(80,80,80))

add_text(sl, Inches(0.35), CONTENT_Y + Inches(2.5), Inches(5.6), Inches(0.28),
         "CAH Ward (n=5 000 — complexité O(n²))", size=14, bold=True, color=BLEU)
add_table(sl, Inches(0.35), CONTENT_Y + Inches(2.83),
          Inches(5.6), Inches(1.5),
    ["Cluster", "Profil", "% Fumeurs", "Systolique"],
    [["0", "Non-fumeurs, HTA",      "2%",    "142 mmHg"],
     ["1", "Fumeurs modérés",       "18%",   "115 mmHg"],
     ["2", "Fumeurs intenses, HTA", "100%",  "142 mmHg"]],
    fs=13, cw=[Inches(0.78), Inches(2.17), Inches(1.25), Inches(1.35)])
add_text(sl, Inches(0.35), CONTENT_Y + Inches(4.45), Inches(5.6), Inches(0.3),
         "CAH isole les fumeurs intenses — K-Means les aurait dilués (méthodes complémentaires)",
         size=12, italic=True, color=RGBColor(80,80,80))

add_table(sl, Inches(0.35), CONTENT_Y + Inches(4.9), Inches(5.6), Inches(1.3),
    ["Méthode", "Silhouette", "Davies-Bouldin"],
    [["K-Means", "0.103", "2.655"],
     ["CAH",     "0.105", "2.721"]],
    fs=13, cw=[Inches(1.8), Inches(1.8), Inches(2.0)])

add_img(sl, FIGS+"kmeans_selection_k.png", Inches(6.1), CONTENT_Y + Inches(0.05), w=Inches(3.6))
caption(sl, Inches(6.1), CONTENT_Y + Inches(3.15), Inches(3.6),
        "Sélection k — k=3 maximise la silhouette")
add_img(sl, FIGS+"kmeans_pca.png", Inches(9.85), CONTENT_Y + Inches(0.05), w=Inches(3.15))
caption(sl, Inches(9.85), CONTENT_Y + Inches(3.15), Inches(3.15),
        "Clusters K-Means — ACP 2D (24.3% variance)")
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "11", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 12 — INTERPRÉTABILITÉ
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Interprétabilité — 5 méthodes comparées")

add_text(sl, Inches(0.35), CONTENT_Y + Inches(0.1), W - Inches(0.7), Inches(0.3),
         "Méthodes : coefficients LR & Ridge  |  Feature Importance Gini (RF & XGBoost)  |  SHAP (XGBoost)",
         size=13, color=NOIR)

add_table(sl, Inches(0.35), CONTENT_Y + Inches(0.5),
          Inches(5.6), Inches(1.85),
    ["Variable", "LR", "RF", "XGB", "SHAP", "Ridge", "Score"],
    [["systolic_bp",    "✓", "✓", "✓", "✓", "✓", "5/5"],
     ["daily_steps",    "✓", "✓", "–", "✓", "–", "3/5"],
     ["family_history", "✓", "–", "✓", "–", "✓", "3/5"],
     ["bmi",            "–", "✓", "–", "✓", "–", "2/5"]],
    fs=13, cw=[Inches(1.8), Inches(0.55), Inches(0.55), Inches(0.55), Inches(0.7), Inches(0.7), Inches(0.7)])

beamer_block(sl, Inches(0.35), CONTENT_Y + Inches(2.5),
             Inches(5.6), Inches(2.6),
    "Analyse critique",
    ["Les 5 méthodes ne convergent pas → mesurent du bruit aléatoire",
     "systolic_bp partout : grande variance, pas signal causal",
     "Divergence = résultat analytique en soi (Rudin 2019)",
     "Sur données réelles : cholestérol, tabagisme, IMC domineraient"])

add_img(sl, FIGS+"rf_feature_importance.png", Inches(6.1), CONTENT_Y + Inches(0.05), w=Inches(3.5))
caption(sl, Inches(6.1), CONTENT_Y + Inches(2.9), Inches(3.5),
        "Feature importance Gini — Random Forest")
add_img(sl, FIGS+"shap_summary.png", Inches(9.8), CONTENT_Y + Inches(0.05), w=Inches(3.2))
caption(sl, Inches(9.8), CONTENT_Y + Inches(2.9), Inches(3.2),
        "SHAP summary plot — XGBoost")
add_img(sl, FIGS+"lr_coefficients.png", Inches(6.1), CONTENT_Y + Inches(3.2), w=Inches(6.9))
caption(sl, Inches(6.1), CONTENT_Y + Inches(5.45), Inches(6.9),
        "Coefficients Logistic Regression (variables standardisées)")
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "12", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 13 — BILAN
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Bilan — Résultat central")

add_rect(sl, Inches(0.35), CONTENT_Y + Inches(0.1), W - Inches(0.7), Inches(0.9), GRIS_CLAIR)
add_rect(sl, Inches(0.35), CONTENT_Y + Inches(0.1), Pt(4), Inches(0.9), ROUGE_AV)
add_text(sl, Inches(0.6), CONTENT_Y + Inches(0.14), W - Inches(1.0), Inches(0.25),
         "Limite fondamentale du dataset synthétique", size=13, bold=True, color=ROUGE_AV)
add_text(sl, Inches(0.6), CONTENT_Y + Inches(0.42), W - Inches(1.0), Inches(0.5),
         "Les variables ont été générées indépendamment — sans relation causale encodée. "
         "Aucune technique ne peut extraire un signal qui n'existe pas.",
         size=14, color=NOIR)

y4 = CONTENT_Y + Inches(1.15)
bullet_list(sl, Inches(0.35), y4, Inches(5.9), Inches(2.5),
    ["Corrélations < 0.01 dès l'EDA le prédisait",
     "Lasso : 16/16 coefficients = 0",
     "R² ≈ 0 et AUC ≈ 0.5 pour tous les modèles",
     "Data leakage détecté et corrigé (R² 0.675 → 0)",
     "Interprétabilité divergente (5 méthodes)"],
    title="Preuves convergentes", t_color=ROUGE_AV, size=14)

bullet_list(sl, Inches(6.6), y4, Inches(6.3), Inches(2.5),
    ["Pipeline complet et rigoureux (CRISP-DM)",
     "Détection du data leakage — piège rare en cours théoriques",
     "Validation croisée stratifiée systématique",
     "Deux stratégies d'imbalance comparées",
     "Directement transposable à des données réelles (Framingham, NHANES)"],
    title="Ce n'est pas un échec — valeur méthodologique", size=14)

add_rect(sl, Inches(0.35), CONTENT_Y + Inches(3.9), W - Inches(0.7), Inches(1.0), BLEU)
add_text(sl, Inches(0.55), CONTENT_Y + Inches(4.1), W - Inches(0.9), Inches(0.7),
         "« Savoir diagnostiquer pourquoi un modèle échoue "
         "est aussi précieux que d'en construire un qui réussit. »",
         size=17, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "13", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 14 — PERSPECTIVES
# ══════════════════════════════════════════════════════════════
sl = new_slide()
title_bar(sl, "Perspectives")

persp = [
    ("1", "Changer de dataset",
     "Cohortes Framingham ou NHANES — vraies relations biologiques, même pipeline"),
    ("2", "Améliorer la classification",
     "Stacking, réseaux de neurones (Deep Learning M2) — SMOTE déjà intégré"),
    ("3", "Déploiement",
     "API REST (FastAPI) : profil individuel → prédiction + valeurs SHAP"),
    ("4", "Interprétabilité avancée",
     "Partial Dependence Plots (PDP), ICE — au-delà de SHAP agrégé"),
    ("5", "Data drift",
     "Suivi AUC sur nouvelles cohortes, réentraînement périodique"),
    ("6", "Master 2",
     "Causalité, fairness algorithmique, apprentissage fédéré"),
]
for i, (num, titre, desc) in enumerate(persp):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.5)
    y = CONTENT_Y + Inches(0.12) + row * Inches(1.7)
    add_rect(sl, x, y, Inches(6.1), Inches(1.5), GRIS_CLAIR)
    add_rect(sl, x, y, Pt(4), Inches(1.5), BLEU_ACC)
    add_rect(sl, x, y, Inches(0.45), Inches(1.5), GRIS_CLAIR)  # fond pastille
    add_rect(sl, x, y, Inches(0.45), Inches(1.5), GRIS_CLAIR)
    # numéro
    add_text(sl, x + Inches(0.04), y + Inches(0.4), Inches(0.4), Inches(0.55),
             num, size=22, bold=True, color=BLEU_ACC, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.52), y + Inches(0.1), Inches(5.45), Inches(0.38),
             titre, size=14, bold=True, color=BLEU)
    add_text(sl, x + Inches(0.52), y + Inches(0.52), Inches(5.45), Inches(0.85),
             desc, size=13, color=NOIR)

add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "14", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# DIAPO 15 — QUESTIONS
# ══════════════════════════════════════════════════════════════
sl = new_slide()
add_rect(sl, 0, Inches(2.9), W, Pt(2), BLEU_ACC)
add_rect(sl, 0, Inches(4.9), W, Pt(2), BLEU_ACC)
add_text(sl, Inches(0.8), Inches(3.1), W - Inches(1.6), Inches(1.1),
         "Merci pour votre attention",
         size=38, bold=True, color=BLEU, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(4.1), W - Inches(1.6), Inches(0.65),
         "Questions ?",
         size=26, color=BLEU_ACC, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(5.2), W - Inches(1.6), Inches(0.35),
         "Ardacham Mahamat Teguene  —  M1 MIASHS  —  Université de Montpellier Paul Valéry  —  Juin 2026",
         size=13, color=RGBColor(120,120,120), align=PP_ALIGN.CENTER)
add_text(sl, W - Inches(0.7), H - Inches(0.35), Inches(0.55), Inches(0.28),
         "15", size=10, color=RGBColor(140,140,140), align=PP_ALIGN.RIGHT)

# ── Sauvegarde ───────────────────────────────────────────────
OUT = ("/Users/ardachammahamatteguene/Desktop/"
       "Bureau - MacBook Pro de Ardacham/Memoire/presentation.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
